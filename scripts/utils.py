"""Shared utilities: PDF/PPTX extraction, Claude client, file I/O.

Auth strategy (in priority order):
  1. ANTHROPIC_API_KEY set  → use Python SDK directly (fast, no subprocess overhead)
  2. No API key             → use `claude -p` CLI subprocess (uses your Claude.ai subscription)
"""

import os
import sys
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Optional

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt
from dotenv import load_dotenv

load_dotenv()

ROOT = Path(__file__).parent.parent
INPUT = ROOT / "input"
OUTPUT = ROOT / "output"


# ---------------------------------------------------------------------------
# Auth detection
# ---------------------------------------------------------------------------

def _use_sdk() -> bool:
    """Return True if we should use the Python SDK (API key available)."""
    return bool(os.environ.get("ANTHROPIC_API_KEY"))


def get_client():
    """Return an anthropic.Anthropic client if API key is set, else None (CLI mode)."""
    if _use_sdk():
        import anthropic
        return anthropic.Anthropic(api_key=os.environ["ANTHROPIC_API_KEY"])
    return None  # signals CLI mode to ask_claude*


# ---------------------------------------------------------------------------
# Claude call — SDK path
# ---------------------------------------------------------------------------

def _ask_sdk(client, system: str, user: str, model: str, max_tokens: int) -> str:
    response = client.messages.create(
        model=model,
        max_tokens=max_tokens,
        system=system,
        messages=[{"role": "user", "content": user}],
    )
    return response.content[0].text


# ---------------------------------------------------------------------------
# Claude call — CLI path (uses Claude.ai subscription)
# ---------------------------------------------------------------------------

def _find_claude_exe() -> str:
    """Locate the claude CLI executable, including Windows npm install locations."""
    # 1. Already on PATH (works in most shells)
    found = shutil.which("claude") or shutil.which("claude.cmd")
    if found:
        return found
    # 2. Windows npm global install (most common for Claude Code)
    appdata = os.environ.get("APPDATA", "")
    for candidate in [
        Path(appdata) / "npm" / "claude.cmd",
        Path(appdata) / "npm" / "claude",
        Path(os.path.expanduser("~")) / "AppData" / "Roaming" / "npm" / "claude.cmd",
    ]:
        if candidate.exists():
            return str(candidate)
    raise FileNotFoundError(
        "Could not find the `claude` CLI. "
        "Make sure Claude Code is installed and you have run it at least once."
    )


def _ask_cli(system: str, user: str) -> str:
    """
    Call `claude -p` via subprocess using the logged-in Claude.ai subscription.
    Pipes the combined system+user prompt via stdin to avoid shell-escaping issues.
    """
    full_prompt = f"{system}\n\n---\n\n{user}"
    claude_exe = _find_claude_exe()

    result = subprocess.run(
        [claude_exe, "-p", "--output-format", "text"],
        input=full_prompt,
        capture_output=True,
        text=True,
        encoding="utf-8",
        timeout=None,  # no cap — large prompts (thesis PDF) can take 10+ min via CLI
    )
    if result.returncode != 0:
        err = result.stderr.strip() or "(no stderr)"
        raise RuntimeError(
            f"`claude -p` exited with code {result.returncode}.\n"
            f"stderr: {err}\n"
            "Make sure you are logged in: run `claude` once interactively first."
        )
    return result.stdout.strip()


# ---------------------------------------------------------------------------
# Public API — used by all agents
# ---------------------------------------------------------------------------

def ask_claude(
    client,  # anthropic.Anthropic or None
    system: str,
    user: str,
    model: str = "claude-opus-4-6",
    max_tokens: int = 8192,
) -> str:
    """Single-turn Claude call. Works with SDK client or CLI (client=None)."""
    if client is not None:
        return _ask_sdk(client, system, user, model, max_tokens)
    else:
        print("  (using claude CLI / Claude.ai subscription — no API key needed)")
        return _ask_cli(system, user)


def ask_claude_long(
    client,
    system: str,
    user: str,
    model: str = "claude-opus-4-6",
    max_tokens: int = 16000,
) -> str:
    """Extended output call for long analysis tasks."""
    return ask_claude(client, system, user, model=model, max_tokens=max_tokens)


# ---------------------------------------------------------------------------
# PDF / PPTX extraction
# ---------------------------------------------------------------------------

def extract_pdf_text(pdf_path: Path, max_chars: int = 400_000) -> str:
    """Extract text from PDF, truncating to max_chars."""
    doc = fitz.open(str(pdf_path))
    pages = []
    total = 0
    for page_num, page in enumerate(doc):
        text = page.get_text()
        pages.append(f"[Page {page_num + 1}]\n{text}")
        total += len(text)
        if total >= max_chars:
            pages.append("\n[... truncated for length ...]")
            break
    return "\n\n".join(pages)


def extract_pdf_by_pages(pdf_path: Path) -> list[dict]:
    """Extract PDF as list of {page, text} dicts."""
    doc = fitz.open(str(pdf_path))
    return [{"page": i + 1, "text": page.get_text()} for i, page in enumerate(doc)]


def extract_pptx_structured(pptx_path: Path) -> list[dict]:
    """Extract PPTX as list of {slide_num, title, body, notes} dicts."""
    prs = Presentation(str(pptx_path))
    slides = []
    for i, slide in enumerate(prs.slides):
        title = ""
        body_parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                if shape.placeholder_format and shape.placeholder_format.idx == 0:
                    title = shape.text_frame.text.strip()
                    continue
            except ValueError:
                pass
            text = shape.text_frame.text.strip()
            if text:
                body_parts.append(text)

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append({
            "slide_num": i + 1,
            "title": title,
            "body": "\n".join(body_parts),
            "notes": notes,
        })
    return slides


def slides_to_text(slides: list[dict], include_notes: bool = True) -> str:
    """Convert structured slides to a readable text block."""
    lines = []
    for s in slides:
        lines.append(f"\n--- Slide {s['slide_num']} ---")
        if s["title"]:
            lines.append(f"TITLE: {s['title']}")
        if s["body"]:
            lines.append(s["body"])
        if include_notes and s["notes"]:
            lines.append(f"[NOTES]: {s['notes']}")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# File I/O
# ---------------------------------------------------------------------------

def save_output(filename: str, content: str, subfolder: str = "") -> Path:
    """Save content to output directory."""
    folder = OUTPUT / subfolder if subfolder else OUTPUT
    folder.mkdir(parents=True, exist_ok=True)
    path = folder / filename
    path.write_text(content, encoding="utf-8")
    print(f"  Saved: {path.relative_to(ROOT)}")
    return path


def load_output(filename: str, subfolder: str = "") -> Optional[str]:
    """Load previously generated output if it exists."""
    folder = OUTPUT / subfolder if subfolder else OUTPUT
    path = folder / filename
    if path.exists():
        return path.read_text(encoding="utf-8")
    return None
